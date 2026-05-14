import base64
import html
import io
import os
import re
import tempfile

import streamlit as st
from fpdf import FPDF
from pptx import Presentation
import google.generativeai as genai

try:
    from dotenv import load_dotenv

    load_dotenv()
except ImportError:
    pass

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from PIL import Image
except ImportError:
    Image = None

# --- AI YAPILANDIRMASI ---
# Anahtar repoda tutulmaz. Siralama: ortam GEMINI_API_KEY / GOOGLE_API_KEY,
# sonra Streamlit secrets. Istege bagli: GEMINI_MODEL (ornegin gemini-3.1-flash-lite).
# Varsayilan: gemini-3.1-flash-lite-preview (cogu ucretsiz planda RPD 500, RPM 15 — Streamlit icin uygun).
# Daha agir kalite: GEMINI_MODEL=gemini-2.5-flash (RPD genelde 20; dusuk trafikte iyi).
# Sablon: .streamlit/secrets.toml.example ve .env.example
def _gemini_api_key() -> str:
    k = (os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY") or "").strip()
    if k:
        return k
    try:
        sec = st.secrets
        for name in ("GEMINI_API_KEY", "GOOGLE_API_KEY"):
            if name in sec:
                v = str(sec[name]).strip()
                if v:
                    return v
    except Exception:
        pass
    return ""


def _gemini_model_name() -> str:
    m = (os.getenv("GEMINI_MODEL") or "").strip()
    if m:
        return m
    try:
        sec = st.secrets
        if "GEMINI_MODEL" in sec:
            return str(sec["GEMINI_MODEL"]).strip()
    except Exception:
        pass
    return "gemini-3.1-flash-lite-preview"


_GEMINI_KEY = _gemini_api_key()
if _GEMINI_KEY:
    genai.configure(api_key=_GEMINI_KEY)
    ai_model = genai.GenerativeModel(
        model_name=_gemini_model_name(),
        system_instruction=(
            "Sen Bridge-AI asistanısın. Merve Yılmaz tarafından geliştirildin. "
            "Analizlerinde profesyonel bir mühendis ve samimi bir öğretmen gibi davran."
        ),
    )
else:
    ai_model = None


def get_ai_response(prompt: str, image_file=None) -> str:
    if ai_model is None:
        return (
            "Gemini API anahtarı yapılandırılmadı. GitHub’a anahtar koymayın. "
            "Yerelde: ortam değişkeni GEMINI_API_KEY veya proje kökünde "
            ".streamlit/secrets.toml (gitignore’da). Streamlit Cloud: "
            "App settings → Secrets."
        )
    try:
        if image_file and Image is not None:
            try:
                image_file.seek(0)
            except (AttributeError, OSError):
                pass
            img = Image.open(image_file)
            response = ai_model.generate_content([prompt, img])
        else:
            response = ai_model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"AI Hatası: {str(e)}"


def _is_likely_exam_question_line(line: str) -> bool:
    """Giris metnini ele; numarali veya soru isareti iceren satirlari soru say."""
    t = line.strip()
    if len(t) < 12:
        return False
    tl = t.lower()
    hints = (
        "merhaba",
        "bridge-ai",
        "merve y",
        "merve yılmaz",
        "analiz ettim",
        "hazırladım",
        "hazirladim",
        "mühendis titiz",
        "muhendis titiz",
        "öğretmen şefkat",
        "ogretmen sefkat",
        "geliştirilen",
        "gelistirilen",
        "asistanın",
        "asistanin",
        "çalışma sorular",
        "calisma sorular",
        "dayalı çalışma",
        "dayali calisma",
        "materyali senin için",
        "keyifli ama dikkat",
    )
    if any(h in tl for h in hints):
        return False
    if tl.startswith("baslik:"):
        return False
    if re.match(r"^\d+[\).]\s*\S", t):
        return True
    if t.rstrip().endswith("?"):
        return True
    if "?" in t and len(t) < 500:
        return True
    return False


def _parse_exam_ai_response(res: str, question_count: int) -> tuple[str, list[str]]:
    """BASLIK: (herhangi bir satir), CEVAP ANAHTARI oncesi sorular; giris metni elenir."""
    m_cut = re.search(
        r"(?im)^[\s#*_-]*(?:CEVAP\s*ANAHTARI|CEVAP\s*ANAHTAR|ANSWER\s*KEY)\b",
        res,
    )
    work = res[: m_cut.start()].strip() if m_cut else res.strip()

    topic = "Ders Sınavı"
    work_lines = work.splitlines()
    body_lines: list[str] = []
    baslik_re = re.compile(r"^\s*BASLIK\s*:\s*(.+)$", re.IGNORECASE)
    for raw in work_lines:
        ln = raw.strip()
        m = baslik_re.match(ln)
        if m:
            cand = m.group(1).strip()
            if cand:
                topic = cand
            continue
        body_lines.append(raw)

    lines = [ln.strip().lstrip("-*• ").strip() for ln in body_lines if ln.strip()]
    qs: list[str] = []
    for ln in lines:
        if _is_likely_exam_question_line(ln):
            qs.append(ln)
            if len(qs) >= question_count:
                break

    if len(qs) < question_count:
        for ln in lines:
            if ln in qs:
                continue
            if len(ln) > 35 and not ln.lower().startswith("baslik"):
                qs.append(ln)
            if len(qs) >= question_count:
                break

    return topic, qs[:question_count]


def _fpdf_latin1_safe(text: str) -> str:
    """FPDF varsayilan fontlari Latin-1 ile calisir; em dash vb. UnicodeEncodeError verir."""
    if not isinstance(text, str):
        text = str(text)
    for old, new in (
        ("\u2014", "-"),  # em dash
        ("\u2013", "-"),  # en dash
        ("\u2018", "'"),
        ("\u2019", "'"),
        ("\u201c", '"'),
        ("\u201d", '"'),
        ("\u2026", "..."),
        ("\u00a0", " "),
    ):
        text = text.replace(old, new)
    return text.encode("latin-1", errors="replace").decode("latin-1")


def create_general_pdf(title, content_list):
    title = _fpdf_latin1_safe(title)
    safe_items = [_fpdf_latin1_safe(item) for item in content_list]
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(190, 10, title, ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("Arial", "", 12)
    for item in safe_items:
        pdf.multi_cell(190, 10, f"- {item}")
        pdf.ln(2)
    return bytes(pdf.output())

def create_presentation_pptx(title, slide_points):
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Bridge-AI Sunum Ciktisi"

    content_layout = prs.slide_layouts[1]
    for point in slide_points:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = "Sunum Icerigi"
        slide.placeholders[1].text = point

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def create_office_excel_bytes():
    if pd is None:
        return None
    df = pd.DataFrame(
        {
            "Kalem": ["Demirbas", "Sarf Malzeme", "Lojistik", "Toplam"],
            "Adet": [3, 42, 8, 53],
            "Birim": ["Adet", "Adet", "Sefer", "-"],
            "Not": ["Ofis raporu ornek", "Stok", "Sevkiyat", "Ozet"],
        }
    )
    buf = io.BytesIO()
    for eng in ("xlsxwriter", "openpyxl"):
        try:
            buf.seek(0)
            buf.truncate(0)
            with pd.ExcelWriter(buf, engine=eng) as writer:
                df.to_excel(writer, index=False, sheet_name="OfisRaporu")
            buf.seek(0)
            return buf.getvalue()
        except Exception:
            continue
    return None


def create_office_scan_pdf(title: str, filename: str, img_bytes=None):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    safe_title = _fpdf_latin1_safe(title)
    pdf.cell(190, 10, safe_title, ln=True, align="C")
    pdf.set_font("Arial", "", 11)
    pdf.multi_cell(190, 8, _fpdf_latin1_safe(f"Taranan / yuklenen dosya: {filename}"))
    pdf.ln(4)
    if img_bytes and Image is not None:
        try:
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            Image.open(io.BytesIO(img_bytes)).convert("RGB").save(tmp.name, format="PNG")
            pdf.image(tmp.name, x=15, w=180)
            try:
                os.unlink(tmp.name)
            except OSError:
                pass
        except Exception:
            pdf.multi_cell(190, 8, _fpdf_latin1_safe("(Gorsel PDFye gomulemedi; dosya adi raporda yer alir.)"))
    else:
        pdf.multi_cell(190, 8, _fpdf_latin1_safe("Gorsel bulunamadi veya PIL yuklu degil; metin ozeti uretildi."))
    return bytes(pdf.output())


def office_slide_points(filename: str):
    return [
        f"Kaynak belge: {filename}",
        "Ozet: Ofis sureclerinde dijitallesme ve raporlama",
        "Onerilen aksiyonlar: kontrol listesi, onay akisi, arsivleme",
        "Riskler: veri kalitesi, zamanlama, uyumluluk",
        "Sonraki adim: paydaslarla gozden gecirme ve yayim",
    ]


def notes_col_html():
    msg = st.session_state.get("ai_msg") or "Bir islem secildiginde burada kisa aciklama gorunur."
    return (
        '<div class="preview-notes-col"><strong>Notlar</strong>'
        '<p style="margin:0;line-height:1.55;">'
        + html.escape(msg)
        + "</p></div>"
    )


def render_preview(placeholder, body_inner_html: str):
    placeholder.markdown(
        '<div class="preview-container preview-flex"><div class="preview-body-col">'
        + body_inner_html
        + "</div>"
        + notes_col_html()
        + "</div>",
        unsafe_allow_html=True,
    )


def go_rerun(celebrate: bool = False):
    if celebrate:
        st.session_state["celebrate"] = True
    st.rerun()


def _optional_app_background_css() -> str:
    """Sabit arka plan: BRIDGE_BG_IMAGE tam yolu veya script yaninda assets/background.jpg|jpeg|png|webp."""
    path = (os.getenv("BRIDGE_BG_IMAGE") or "").strip()
    if not path:
        base = os.path.dirname(os.path.abspath(__file__))
        for name in (
            "background.jpg",
            "background.jpeg",
            "background.png",
            "background.webp",
        ):
            cand = os.path.join(base, "assets", name)
            if os.path.isfile(cand):
                path = cand
                break

    if path.startswith("http://") or path.startswith("https://"):
        return (
            ".stApp {"
            "background-image: url('" + path + "');"
            "background-size: cover;"
            "background-position: center center;"
            "background-attachment: fixed;"
            "background-repeat: no-repeat;"
            "}"
            ".stApp > .main { background-color: transparent !important; }"
            ".main .block-container {"
            "background-color: rgba(248, 249, 250, 0.93);"
            "border-radius: 12px;"
            "padding-top: 1rem;"
            "}"
        )

    if not path or not os.path.isfile(path):
        return ""

    ext = os.path.splitext(path)[1].lower()
    mime = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
        ".gif": "image/gif",
    }.get(ext, "image/jpeg")
    try:
        with open(path, "rb") as f:
            b64 = base64.standard_b64encode(f.read()).decode("ascii")
    except OSError:
        return ""
    return (
        ".stApp {"
        "background-image: url('data:" + mime + ";base64," + b64 + "');"
        "background-size: cover;"
        "background-position: center center;"
        "background-attachment: fixed;"
        "background-repeat: no-repeat;"
        "}"
        ".stApp > .main { background-color: transparent !important; }"
        ".main .block-container {"
        "background-color: rgba(248, 249, 250, 0.93);"
        "border-radius: 12px;"
        "padding-top: 1rem;"
        "}"
    )

    if not path or not os.path.isfile(path):
        return ""

    ext = os.path.splitext(path)[1].lower()
    mime = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
        ".gif": "image/gif",
    }.get(ext, "image/jpeg")
    try:
        with open(path, "rb") as f:
            b64 = base64.standard_b64encode(f.read()).decode("ascii")
    except OSError:
        return ""
    return f"""
    .stApp {{
        background-image: url("data:{mime};base64,{b64}");
        background-size: cover;
        background-position: center center;
        background-attachment: fixed;
        background-repeat: no-repeat;
    }}
    .stApp > .main {{ background-color: transparent !important; }}
    .main .block-container {{
        background-color: rgba(248, 249, 250, 0.93);
        border-radius: 12px;
        padding-top: 1rem;
    }}
    """
    .stApp > .main {{ background-color: transparent !important; }}
    .main .block-container {{
        background-color: rgba(248, 249, 250, 0.93);
        border-radius: 12px;
        padding-top: 1rem;
    }}
    """

if not path or not os.path.isfile(path):
    return ""
    ext = os.path.splitext(path)[1].lower()
    mime = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
        ".gif": "image/gif",
    }.get(ext, "image/jpeg")
    try:
        with open(path, "rb") as f:
            b64 = base64.standard_b64encode(f.read()).decode("ascii")
    except OSError:
        return ""
    return f"""
    .stApp {{
        background-image: url("data:{mime};base64,{b64}");
        background-size: cover;
        background-position: center center;
        background-attachment: fixed;
        background-repeat: no-repeat;
    }}
    .main {{ background-color: transparent !important; }}
    .main .block-container {{
        background-color: rgba(248, 249, 250, 0.93);
        border-radius: 12px;
        padding-top: 1rem;
    }}
    """


# --- 2. KURUMSAL TASARIM VE CSS ---
st.set_page_config(page_title="BRIDGE-AI® Enterprise PRO", layout="wide")

st.markdown(
    """
    <div id="bridge-boot" style="position:fixed;inset:0;background:rgba(248,250,252,0.94);z-index:99999;
    display:flex;flex-direction:column;align-items:center;justify-content:center;animation:bootFade 2.2s ease forwards;">
      <div style="font-size:3.5rem;animation:botBounce 0.75s ease-in-out infinite alternate;">🤖</div>
      <p style="font-weight:600;color:#0f172a;">Bridge-AI yukleniyor...</p>
    </div>
    <style>
    @keyframes bootFade { 0%{opacity:1} 75%{opacity:1} 100%{opacity:0;visibility:hidden;pointer-events:none;} }
    @keyframes botBounce { from { transform: translateY(0); } to { transform: translateY(-10px); } }
    </style>
    """,
    unsafe_allow_html=True,
)

_app_background_css = _optional_app_background_css()
_main_background_rule = (
    "    .main { background-color: #f8f9fa; }\n"
    if not _app_background_css
    else ""
)
st.markdown(
    """
    <style>
"""
    + _app_background_css
    + _main_background_rule
    + """
    [data-testid="stSidebar"] { background-color: #001f3f !important; }
    [data-testid="stSidebar"] * { color: white !important; font-weight: 600; }
    
    /* KPI KARTLARI (İşaretlediğin Tasarım) */
    .kpi-container {
        display: flex;
        justify-content: space-between;
        gap: 10px;
        margin-bottom: 20px;
    }
    .kpi-card {
        background-color: white;
        padding: 15px 25px;
        border-radius: 12px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        flex: 1;
        border: 1px solid #e2e8f0;
        display: flex;
        flex-direction: column;
        align-items: flex-start;
    }
    .kpi-card span { font-size: 0.8rem; color: #64748b; text-transform: uppercase; letter-spacing: 1px; }
    .kpi-card strong { font-size: 1.4rem; color: #0f172a; margin-top: 5px; }
    
    /* ÖNİZLEME KUTUSU (kesik cizgi tum satiri kaplar; notlar ic kisimda yan bant) */
    .preview-container {
        border: 2px dashed #10B981 !important;
        border-radius: 15px;
        background-color: #ffffff;
        min-height: 450px;
        padding: 22px;
        margin-bottom: 20px;
    }
    .preview-container.preview-flex {
        display: flex;
        flex-direction: row;
        align-items: flex-start;
        gap: 18px;
        flex-wrap: wrap;
    }
    .preview-body-col {
        flex: 1.55;
        min-width: 260px;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: flex-start;
    }
    .preview-notes-col {
        flex: 1;
        min-width: 220px;
        background: transparent;
        border-left: 3px dashed #10B981;
        padding: 8px 8px 8px 16px;
        color: #0f172a;
    }

    .paper-mockup {
        background-color: white; border: 1px solid #e2e8f0; padding: 30px;
        box-shadow: 0 10px 25px rgba(0,0,0,0.1); font-family: 'serif'; text-align: left; color: #1e293b;
        width: 100%; max-width: 600px;
    }
    
    .op-card { background-color: white; padding: 25px; border-radius: 12px; border: 1px solid #e2e8f0; }
    </style>
    """,
    unsafe_allow_html=True,
)

# --- 3. SIDEBAR ---
with st.sidebar:
    st.title("BRIDGE-AI®")
    module = st.radio("NAVİGASYON", ["Eğitim & Akademi", "Saha (Field)", "Ofis & İdari İşler", "Yönetici Paneli"])
    st.write("👤 Merve Yılmaz Çiftçi")

# --- 4. ÜST KPI PANELİ (İşaretlediğin Görseldeki Tasarım) ---
st.markdown(f"""
    <div class="kpi-container">
        <div class="kpi-card"><span>TOTAL ANALYSES</span><strong>5.331</strong></div>
        <div class="kpi-card"><span>CRITICAL RISKS</span><strong style="color:#ef4444">3 🔴</strong></div>
        <div class="kpi-card"><span>AI SPEED</span><strong style="color:#10B981">4h 4m/s</strong></div>
        <div class="kpi-card"><span>PENDING TASKS</span><strong>2</strong></div>
    </div>
""", unsafe_allow_html=True)

st.header(f"📊 {module} Kontrol Paneli")

with st.expander("Bu arayüzde neler var? (güncel özellikler)", expanded=True):
    st.markdown(
        """
- **Sunum:** kesik çizgili kutuda slayt maddeleri + altta **Sunumu İndir (.pptx)**  
- **Sınav:** soru sayısı + kutuda önizleme + **Sınavı PDF İndir**  
- **Kelime bankası:** özet ve hap bilgiler + **Kelime Özetini İndir**  
- **Saha:** rapor önizleme + indirme | **SOP** taslağı + indirme  
- **Ofis:** OCR + Excel raporu + **ayri** gorsel sunum butonu; indirmeler asagida  
- **Yönetici:** stratejik özet önizleme | risk raporu + indirme  
- **Notlar:** kesik cizgili cercevenin **icinde**, onizleme ile yan yana (ayri beyaz kart yok)  
"""
    )

# --- 5. ANA PANEL VE ÖNİZLEME ---
col_left, col_right = st.columns([2.3, 1])

# Hazır İçerik Havuzu (Egitim: Mars Survival + Conditionals / 11. sinif Ingilizce ile uyumlu)
exam_questions = [
    "Write a paragraph (120-150 words): You are on the Mars Survival Council. Explain three priorities (water, energy, teamwork) and what will happen if the team ignores one priority. Use at least two future-tense sentences.",
    "Second Conditional: Complete the sentences. (a) If we _____ (have) more solar panels, we _____ (save) more energy. (b) If the habitat _____ (be) colder, we _____ (need) extra insulation.",
    "Rewrite using Unless: 'If we don't share the water fairly, the mission will fail.'",
    "Role card speaking (write 6-8 sentences): You must convince the council to approve your Bridge-AI logistics plan. Use one second conditional and one sentence with unless.",
    "Reading-to-writing: Read the scenario title 'Mars Survival Council - Energy Crisis'. Write five bullet decisions your council will take next week. Each bullet must start with 'We will ...'."
]
presentation_points = [
    "Unite basligi: Mars Survival Council - Energy & Cooperation (11. sinif Ingilizce)",
    "Ogrenme ciktilari: Second Conditional / Unless / Future plans ile gerekce sunma",
    "Isinma: Mars ortami kelime bankasi + kisa gorsel taslak (3 madde)",
    "Anlatim: Conditionals ornekleri + sinif aktivitesi (cift calismasi + paylasim)",
    "Degerlendirme: rubrik (dil dogrulugu, baglam, yapilar) + Bridge-AI proje baglantisi (1 slayt)"
]
vocab_summary_points = [
    "Konu ozeti: 11. sinif Ingilizce Mars Survival Council; enerji krizi ve is birligi ile karar alma baglami.",
    "Hap bilgi: Second Conditional = If + past simple, would + fiil (hayali simdiki-gelecek).",
    "Hap bilgi: Unless = 'if ... not' anlami; Unless we save energy, we will fail.",
    "Hap bilgi: Future plan cumleleri: will / be going to + eylem + zaman ifadesi (next week).",
    "Hap bilgi: Paragraf yazarken once priority (3), sonra sonuc (if ignore), en az 2 future tense.",
]
field_5s_body = """<h3>ISG ve 5S Denetim Bulgulari</h3><hr>
<p><b>Genel:</b> Uretim alaninda 5S ve temel ISG kurallarina uyumsuzluk tespit edildi.</p>
<ul style="text-align:left;line-height:1.55;">
<li><b>Seiri (Ayristir):</b> Calisma yolunda kullanilmayan palet ve ambalaj kalintilari; gereksiz malzeme trafo ve dolap onunu tikiyor.</li>
<li><b>Seiton (Duzenle):</b> Raf etiketleri eksik; aletler sabit yerde degil; acil cikis yonu gecici olarak engellenmis.</li>
<li><b>Seiso (Temizle):</b> Talaş ve yag lekeleri zeminde birikim; kayma riski ve yangin surukleyici tehlike.</li>
<li><b>Seiketsu (Standartlastir):</b> Temizlik ve kontrol listesi uygulanmiyor; ayni hatada tekrarlayan sapmalar.</li>
<li><b>Sitsuke (Disiplin):</b> KKD (baret, is ayakkabisi) kismi kullanim; uyari bantlari solmus veya kopmus.</li>
</ul>
<p><b>ISG oncelikleri:</b> Kaynak yaralanmasi, dusme-carpma, elektrik alani yaklasim, ergonomik zorlanma.</p>"""
sop_points = [
    "Amac ve kapsam tanimi",
    "Guvenlik adimlari ve kontrol noktasi",
    "Standart uygulama akisi",
    "Sorumluluk matrisi",
    "Sapma durumunda aksiyon plani"
]
risk_points = [
    "Kritik risk: is guvenligi ihlali olasiligi",
    "Etki seviyesi: yuksek",
    "Onerilen aksiyon: denetim sikligini artir",
    "Sorumlu ekip: saha yonetimi",
    "Takip KPI: olay frekansi ve kapanis suresi"
]


def _office_is_image(name: str) -> bool:
    return str(name).lower().endswith((".png", ".jpg", ".jpeg", ".webp"))


with col_left:
    preview_placeholder = st.empty()

    if "view_state" not in st.session_state:
        render_preview(
            preview_placeholder,
            '<div style="text-align:center;padding:36px 12px;color:#64748b;">📂 İçerikler burada önizlenecek; aşağıdan modül seçin.</div>',
        )
    else:
        state = st.session_state["view_state"]

        if state == "pdf":
            q_count = st.session_state.get("q_count", 5)
            qs = st.session_state.get("edu_exam_questions")
            display_qs = (qs if qs else exam_questions)[:q_count]
            exam_topic = st.session_state.get("edu_exam_topic") or st.session_state.get("edu_exam_title") or "Ders Sınavı"
            q_html = "".join(
                [f"<li style='margin-bottom:10px;'>{html.escape(q)}</li>" for q in display_qs]
            )
            render_preview(
                preview_placeholder,
                '<div class="paper-mockup">'
                f'<h2 style="text-align:center;margin:0 0 6px 0;color:#0f172a;font-size:1.35rem;">{html.escape(exam_topic)}</h2>'
                '<h3 style="text-align:center;margin:0 0 18px 0;color:#64748b;font-size:1.05rem;font-weight:600;">Sınav soruları</h3>'
                f'<ol style="text-align:left;line-height:1.55;padding-left:1.25rem;">{q_html}</ol>'
                "</div>",
            )
            st.download_button(
                "📥 Sınavı PDF İndir",
                create_general_pdf(exam_topic, display_qs),
                "sinav.pdf",
                "application/pdf",
                use_container_width=True,
            )

        elif state == "slide":
            slide_pts = st.session_state.get("edu_slide_points") or presentation_points
            slide_html = "".join(
                [f"<li style='margin-bottom:10px;'>{html.escape(item)}</li>" for item in slide_pts]
            )
            slide_title = st.session_state.get("edu_slide_title") or "MARS SURVIVAL COUNCIL - UNITE SUNUMU"
            render_preview(
                preview_placeholder,
                f'<div style="background:#001f3f;color:white;padding:40px;border-radius:10px;width:100%;border-left:12px solid #10B981;"><h3>{html.escape(slide_title)}</h3><hr><ol style="padding-left:20px;">{slide_html}</ol></div>',
            )
            st.download_button(
                "📥 Sunumu İndir (.pptx)",
                create_presentation_pptx(slide_title, slide_pts),
                "sunum.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

        elif state == "vocab":
            vocab_pts = st.session_state.get("edu_vocab_points") or vocab_summary_points
            vocab_html = "".join(
                [f"<li style='margin-bottom:10px;'>{html.escape(item)}</li>" for item in vocab_pts]
            )
            render_preview(
                preview_placeholder,
                f'<div class="paper-mockup"><h3>🎓 KELIME BANKASI - OZET VE HAP BILGILER</h3><hr><ol>{vocab_html}</ol></div>',
            )
            st.download_button(
                "📥 Kelime Özetini İndir",
                create_general_pdf("KELIME BANKASI - OZET VE HAP BILGILER", vocab_pts),
                "kelime_ozeti.pdf",
                "application/pdf",
                use_container_width=True,
            )

        elif state == "field_report":
            render_preview(
                preview_placeholder,
                f'<div class="paper-mockup" style="border-top:10px solid #ef4444;">{field_5s_body}</div>',
            )
            field_pdf_lines = [
                "5S: Seiri ihlali - gereksiz malzeme yollarda",
                "5S: Seiton ihlali - etiketsiz raf, alet sabit degil",
                "5S: Seiso ihlali - yag/talas birikimi",
                "5S: Seiketsu ihlali - kontrol listesi uygulanmiyor",
                "5S: Sitsuke ihlali - KKD ve uyari bantlari zayif",
                "ISG: kaynak, dusme, elektrik, ergonomi riskleri",
            ]
            st.download_button(
                "📥 Saha Raporunu İndir",
                create_general_pdf("SAHA 5S ve ISG RAPORU", field_pdf_lines),
                "saha_raporu.pdf",
                "application/pdf",
                use_container_width=True,
            )

        elif state == "office_ocr":
            fn = st.session_state.get("off_last_name", "Belge yok")
            body = (
                f'<div class="paper-mockup"><h3>🧾 OCR VERİ TABLOSU</h3>'
                f"<p>Kaynak dosya: <b>{html.escape(fn)}</b></p>"
                '<table style="width:100%;"><tr><th>Kalem</th><th>Miktar</th><th>Birim</th></tr>'
                "<tr><td>Hammadde A</td><td>500</td><td>Adet</td></tr>"
                "<tr><td>KDV</td><td>20</td><td>%</td></tr>"
                "<tr><td>Toplam</td><td>6000</td><td>TL</td></tr></table>"
                "<p><small>Not: Gercek OCR icin API entegrasyonu eklenebilir; simdi ornek tablo.</small></p></div>"
            )
            render_preview(preview_placeholder, body)
            st.download_button(
                "📥 OCR Satirlarini PDF Indir",
                create_general_pdf("OCR DATA SHEET", ["Hammadde A - 500", "KDV - 20%", "Toplam - 6000 TL"]),
                "ofis_ocr_veri.pdf",
                "application/pdf",
                use_container_width=True,
            )
            img_b = st.session_state.get("off_last_bytes")
            if _office_is_image(fn) and img_b:
                st.download_button(
                    "📥 Taranan Gorseli PDF Indir",
                    create_office_scan_pdf("OFIS TARAMA RAPORU", fn, img_b),
                    "ofis_tarama.pdf",
                    "application/pdf",
                    use_container_width=True,
                )

        elif state == "office_excel":
            render_preview(
                preview_placeholder,
                '<div class="paper-mockup"><h3>📊 EXCEL RAPORU</h3><hr>'
                "<p>Ofis kalemleri tabloya aktarildi. Asagidan <b>.xlsx</b> dosyasini indirebilirsiniz.</p></div>",
            )
            xbytes = create_office_excel_bytes()
            if xbytes:
                st.download_button(
                    "📥 Excel Raporunu Indir (.xlsx)",
                    xbytes,
                    "ofis_raporu.xlsx",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
            else:
                st.warning("Excel icin `pandas` ve `openpyxl` veya `xlsxwriter` gerekir: pip install pandas openpyxl")
            st.download_button(
                "📥 Excel Ozetini PDF Indir",
                create_general_pdf("EXCEL OZET", ["Rapor satirlari Excelde", "Detay xlsx dosyasinda"]),
                "excel_ozet.pdf",
                "application/pdf",
                use_container_width=True,
            )

        elif state == "office_slide":
            fn = st.session_state.get("off_last_name", "belge")
            pts = office_slide_points(fn)
            slide_html = "".join([f"<li style='margin-bottom:10px;'>{html.escape(p)}</li>" for p in pts])
            render_preview(
                preview_placeholder,
                f'<div style="background:#0f172a;color:white;padding:36px;border-radius:10px;border-left:10px solid #38bdf8;">'
                f"<h3>OFIS SUNUM TASLAK</h3><p>Kaynak: {html.escape(fn)}</p><hr><ol>{slide_html}</ol></div>",
            )
            st.download_button(
                "📥 Ofis Sunumunu Indir (.pptx)",
                create_presentation_pptx("OFIS SUNUM TASLAK", pts),
                "ofis_sunum.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
            img_b = st.session_state.get("off_last_bytes")
            if _office_is_image(fn) and img_b:
                st.download_button(
                    "📥 Gorselli PDF Sunum Indir",
                    create_office_scan_pdf("OFIS GORSEL SUNUM PDF", fn, img_b),
                    "ofis_gorsel_sunum.pdf",
                    "application/pdf",
                    use_container_width=True,
                )

        elif state == "exec_summary":
            render_preview(
                preview_placeholder,
                '<div class="paper-mockup" style="border-left:10px solid #f97316;"><h3>📊 YÖNETİCİ ÖZETİ</h3><hr>'
                "<p>Stratejik verimlilik metrikleri ve departman bazlı KPI analizleri optimize edilmiştir.</p></div>",
            )

        elif state == "sop":
            sop_pts = st.session_state.get("field_sop_points") or sop_points
            sop_html = "".join([f"<li style='margin-bottom:10px;'>{html.escape(item)}</li>" for item in sop_pts])
            render_preview(
                preview_placeholder,
                f'<div class="paper-mockup"><h3>📝 SOP TASLAGI</h3><hr><ol>{sop_html}</ol></div>',
            )
            st.download_button(
                "📥 SOP Raporunu İndir",
                create_general_pdf("SOP TASLAGI", sop_pts),
                "sop_raporu.pdf",
                "application/pdf",
                use_container_width=True,
            )

        elif state == "risk_report":
            risk_pts = st.session_state.get("exec_risk_points") or risk_points
            risk_html = "".join([f"<li style='margin-bottom:10px;'>{html.escape(item)}</li>" for item in risk_pts])
            render_preview(
                preview_placeholder,
                f'<div class="paper-mockup" style="border-top:10px solid #ef4444;"><h3>🚨 RISK ANALIZ RAPORU</h3><hr><ol>{risk_html}</ol></div>',
            )
            st.download_button(
                "📥 Risk Raporunu İndir",
                create_general_pdf("RISK ANALIZ RAPORU", risk_pts),
                "risk_analiz_raporu.pdf",
                "application/pdf",
                use_container_width=True,
            )

    # --- İŞLEM KARTI ---
    st.markdown('<div class="op-card">', unsafe_allow_html=True)
    
    if module == "Eğitim & Akademi":
        st.subheader("📚 Eğitim İşlemleri")
        st.file_uploader("Dosya Seç", type=['pdf', 'png', 'jpg'], key="edu_u")
        b1, b2, b3 = st.columns(3)
        if b1.button("✨ Sunum Hazırla", use_container_width=True):
            upl = st.session_state.get("edu_u")
            with st.spinner("AI sunum hazırlıyor..."):
                res = get_ai_response("Bu ders materyalini profesyonel bir sunum taslağına dönüştür. Madde madde yaz.", upl)
            st.session_state['view_state'] = 'slide'
            st.session_state['ai_msg'] = res
            st.session_state["edu_slide_points"] = [line.strip("- ") for line in res.split("\n") if line.strip()][:8]
            st.session_state["edu_slide_title"] = (f"SUNUM — {upl.name}" if upl else "MARS SURVIVAL COUNCIL - UNITE SUNUMU")
            go_rerun(True)
        if b2.button("📄 Sınav Üret", use_container_width=True):
            st.session_state['show_slider'] = True
        if b3.button("🎓 Kelime Bankası", use_container_width=True):
            upl = st.session_state.get("edu_u")
            with st.spinner("AI kelime bankası oluşturuyor..."):
                res = get_ai_response("Bu ders materyalindeki anahtar kelimeleri ve konuyu 11. sınıf seviyesinde özetle. Hap bilgiler ver.", upl)
            st.session_state['view_state'] = 'vocab'
            st.session_state['ai_msg'] = res
            st.session_state["edu_vocab_points"] = [line.strip("- ") for line in res.split("\n") if line.strip()][:6]
            go_rerun(True)
        if st.button("✨ Sunum Hazırla (Yedek)", use_container_width=True, key="sunum_yedek"):
            upl = st.session_state.get("edu_u")
            with st.spinner("AI sunum hazırlıyor..."):
                res = get_ai_response("Bu ders materyalini profesyonel bir sunum taslağına dönüştür. Madde madde yaz.", upl)
            st.session_state['view_state'] = 'slide'
            st.session_state['ai_msg'] = res
            st.session_state["edu_slide_points"] = [line.strip("- ") for line in res.split("\n") if line.strip()][:8]
            st.session_state["edu_slide_title"] = (f"SUNUM — {upl.name}" if upl else "MARS SURVIVAL COUNCIL - UNITE SUNUMU")
            go_rerun(True)
        if st.session_state.get('show_slider'):
            cnt = st.select_slider("Soru Sayısı", options=[5, 10, 15], value=5)
            if st.button("Kutuda Sınavı Oluştur", use_container_width=True):
                upl = st.session_state.get("edu_u")
                base = (
                    f"Bu ders materyaliyle ilgili tam {cnt} adet test sorusu uret. "
                    "Cevap anahtarini EN SONDA ayri bir baslikla yaz: satir basinda tam olarak "
                    "'CEVAP ANAHTARI' yaz, altinda cevaplari ver.\n\n"
                    "KISITLAR: Kendini tanitma, Bridge-AI veya Merve Yilmaz hakkinda paragraf YAZMA; "
                    "giris veya 'merhaba' metni YAZMA. Ilk anlamli satir mutlaka su formatta olsun "
                    "(dosya adi yazma):\n"
                    "BASLIK: <konu / unite, ornek: Wish Clauses / Ingilizce Dilbilgisi>\n\n"
                    "Hemen ardindan (bir bos satir opsiyonel) numarali sorulari yaz; her soru net ve tek baslik altinda okunur olsun."
                )
                if upl:
                    base += (
                        " Yuklenen gorsel veya PDFdeki icerigi incele; sorulari dogrudan bu icerige dayandir."
                    )
                with st.spinner(f"{cnt} soru hazırlanıyor..."):
                    res = get_ai_response(base, upl)
                st.session_state['view_state'] = 'pdf'
                st.session_state['q_count'] = cnt
                st.session_state['ai_msg'] = res
                exam_topic, exam_qs = _parse_exam_ai_response(res, cnt)
                st.session_state["edu_exam_topic"] = exam_topic
                st.session_state["edu_exam_title"] = exam_topic
                st.session_state["edu_exam_questions"] = exam_qs
                go_rerun(True)

    elif module == "Saha (Field)":
        st.subheader("🏭 Saha Denetim & İSG")
        st.file_uploader("Fotoğraf Yükle", type=['png', 'jpg'], key="field_u")
        f1, f2 = st.columns(2)
        if f1.button("🔍 İSG & 5S Analizi", use_container_width=True):
            upl = st.session_state.get("field_u")
            with st.spinner("AI saha analizi yapıyor..."):
                res = get_ai_response("Bu saha fotoğrafını İSG ve 5S açısından denetle. Riskleri ve uyumsuzlukları madde madde listele. Düzeltici faaliyet öner.", upl)
            st.session_state['view_state'] = 'field_report'
            st.session_state['ai_msg'] = res
            go_rerun(True)
        if f2.button("📝 SOP Üretici", use_container_width=True):
            upl = st.session_state.get("field_u")
            with st.spinner("AI SOP hazırlıyor..."):
                res = get_ai_response("Bu saha görseline göre Standart Operasyon Prosedürü (SOP) taslağı oluştur. Adım adım yaz.", upl)
            st.session_state['view_state'] = 'sop'
            st.session_state['ai_msg'] = res
            st.session_state["field_sop_points"] = [line.strip("- ") for line in res.split("\n") if line.strip()][:8]
            go_rerun(True)

    elif module == "Ofis & İdari İşler":
        st.subheader("📂 Ofis Otomasyonu")
        fu_off = st.file_uploader(
            "Fatura PDF, Excel veya fotograf (PNG/JPG)",
            type=["pdf", "xlsx", "png", "jpg", "jpeg"],
            key="off_u",
        )
        o1, o2, o3 = st.columns(3)

        def _store_office_upload(upl):
            if not upl:
                st.session_state["off_last_name"] = "Belge yok"
                st.session_state.pop("off_last_bytes", None)
                return
            st.session_state["off_last_name"] = upl.name
            data = upl.getvalue()
            if len(data) > 5_000_000:
                st.session_state["ai_msg"] = "Dosya 5MB uzeri; daha kucuk bir dosya yukleyin."
                st.session_state.pop("off_last_bytes", None)
            else:
                st.session_state["off_last_bytes"] = data

        if o1.button("🧾 Fatura OCR Analizi", use_container_width=True):
            _store_office_upload(fu_off)
            with st.spinner("AI belgeyi analiz ediyor..."):
                res = get_ai_response("Bu belgedeki tüm verileri (tarih, tutar, vergi no, kalemler) tabloya dök. Satır satır yaz.", fu_off)
            st.session_state["view_state"] = "office_ocr"
            st.session_state["ai_msg"] = res
            go_rerun(True)
        if o2.button("📊 Excel Raporu Uret", use_container_width=True):
            _store_office_upload(fu_off)
            with st.spinner("AI rapor oluşturuyor..."):
                res = get_ai_response("Bu belgedeki kalemleri Excel için tablo formatında listele. Her satırı ayrı yaz.", fu_off)
            st.session_state["view_state"] = "office_excel"
            st.session_state["ai_msg"] = res
            go_rerun(True)
        if o3.button("📽 Gorsel Sunum Uret", use_container_width=True):
            _store_office_upload(fu_off)
            with st.spinner("AI sunum hazırlıyor..."):
                res = get_ai_response("Bu belgeyi yöneticiler için görsel bir sunum taslağına dönüştür. Madde madde yaz.", fu_off)
            st.session_state["view_state"] = "office_slide"
            st.session_state["ai_msg"] = res
            go_rerun(True)

    elif module == "Yönetici Paneli":
        st.subheader("📈 Stratejik Karar Destek")
        y1, y2 = st.columns(2)
        if y1.button("📊 Stratejik Özet Rapor", use_container_width=True):
            with st.spinner("AI stratejik özet hazırlıyor..."):
                res = get_ai_response("Bir üretim şirketi için stratejik özet rapor hazırla. Kritik riskler, stratejik fırsatlar ve karar destek özeti sun.")
            st.session_state['view_state'] = 'exec_summary'
            st.session_state['ai_msg'] = res
            go_rerun(True)
        if y2.button("🚨 Risk Analiz Raporu", use_container_width=True):
            with st.spinner("AI risk analizi yapıyor..."):
                res = get_ai_response("Bir üretim şirketi için detaylı risk analizi raporu hazırla. Her riski etki seviyesi ve önerilen aksiyonla listele.")
            st.session_state['view_state'] = 'risk_report'
            st.session_state['ai_msg'] = res
            st.session_state["exec_risk_points"] = [line.strip("- ") for line in res.split("\n") if line.strip()][:8]
            go_rerun(True)

    st.markdown('</div>', unsafe_allow_html=True)

# --- 6. SAĞ PANEL ---
with col_right:
    st.markdown('<div class="op-card" style="border-left: 5px solid #10B981; min-height:550px;">', unsafe_allow_html=True)
    st.subheader("💡 BRIDGE-AI®")
    st.caption("Notlar, kesik cizgili onizleme kutusunun icinde sag bolumdedir.")
    st.info("Kaynak dosya: `BridgeAI_Streamlit/app.py` — sayfayi yenileyerek son degisiklikleri gor.")
    st.markdown('</div>', unsafe_allow_html=True)

if st.session_state.get("celebrate"):
    st.balloons()
    st.session_state.pop("celebrate", None)
